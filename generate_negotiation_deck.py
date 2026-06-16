"""
generate_negotiation_deck.py
"Difficult Conversations" weekly de-escalation curriculum deck.

Matches the same module shape as generate_praise_deck / generate_wwyd_deck etc:
    generate_region(payload, outdir, region) -> path | None

`payload` is the single-lesson JSON object Claude returns. `region` ("aunz" or
"america") comes from the route and wins over whatever is in the payload.
State (which lesson we're on) lives in Supabase and is handled in Make; this
generator just renders whatever lesson JSON it is handed.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette: deep slate + warm copper ----
DARK   = "1E2A32"
PANEL  = "27353E"
COPPER = "D98E3F"
COPPERD = "B5722D"
CREAM  = "F4F1EC"
CARD   = "FFFFFF"
BODY   = "2A2A2A"
MUTED  = "8A95A0"
STEEL  = "6B7B85"
SAGE   = "6B9B7A"
CARD_LINE = "E5E2DC"
UPCOMING_LIGHT = "E2DDD4"

FONT = "Calibri"
TIER_COUNT = 4
SLIDE_W = 10.0
SLIDE_H = 5.625


def H(hexstr):
    return RGBColor.from_string(hexstr)


def _set_char_spacing(run, pts):
    """Letter spacing in points -> OOXML spc (1/100 pt)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def _no_shadow(shape):
    """Disable inherited default shadow for a clean flat look."""
    shape.shadow.inherit = False


def add_bg(slide, color):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = H(color)
    rect.line.fill.background()
    _no_shadow(rect)
    return rect


def add_round(slide, x, y, w, h, fill_hex, line_hex=None, line_pt=0.0, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = H(fill_hex)
    if line_hex and line_pt:
        shp.line.color.rgb = H(line_hex)
        shp.line.width = Pt(line_pt)
    else:
        shp.line.fill.background()
    _no_shadow(shp)
    return shp


def add_rect(slide, x, y, w, h, fill_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = H(fill_hex)
    shp.line.fill.background()
    _no_shadow(shp)
    return shp


def add_text(slide, x, y, w, h, runs, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, wrap=True):
    """
    runs: a string, or a list of (text, {opts}) tuples for rich runs.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align

    if isinstance(runs, str):
        runs = [(runs, {})]

    for text, opts in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = opts.get("font", FONT)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.color.rgb = H(opts.get("color", color))
        sp = opts.get("spacing", spacing)
        if sp is not None:
            _set_char_spacing(r, sp)
    return tb


def progress_bar(slide, x, y, w, current_tier, lesson_number, total_lessons, on_dark):
    gap = 0.08
    seg_w = (w - gap * (TIER_COUNT - 1)) / TIER_COUNT
    for i in range(TIER_COUNT):
        tier_no = i + 1
        if tier_no < current_tier:
            fill = COPPERD
        elif tier_no == current_tier:
            fill = COPPER
        else:
            fill = PANEL if on_dark else UPCOMING_LIGHT
        add_round(slide, x + i * (seg_w + gap), y, seg_w, 0.14, fill, radius=0.5)
    add_text(slide, x, y + 0.2, w, 0.22,
             "Lesson %s of %s" % (lesson_number, total_lessons),
             9, MUTED, spacing=1)


def bubble(slide, x, y, w, who, line, accent_color):
    h = 0.92
    add_text(slide, x, y, w, 0.2, (who or "").upper(), 8, accent_color,
             bold=True, spacing=1.5)
    add_round(slide, x, y + 0.22, w, h - 0.22, CARD, line_hex=CARD_LINE,
              line_pt=1.0, radius=0.12)
    add_text(slide, x + 0.12, y + 0.3, w - 0.24, h - 0.38, line or "", 10.5,
             BODY, anchor=MSO_ANCHOR.TOP)
    return y + h + 0.12


def derive_tier(lesson_number):
    n = lesson_number
    if n <= 6:
        return 1, "Foundations"
    if n <= 10:
        return 2, "Steering"
    if n <= 15:
        return 3, "Reading"
    return 4, "The Hard Cases"


# ---------- Slides ----------
def cover_slide(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_text(s, 0.5, 0.42, 7, 0.25, "DIFFICULT CONVERSATIONS", 10, COPPER,
             bold=True, spacing=3.5)
    add_text(s, 4.5, 0.42, 5, 0.25,
             "TIER %s OF %s \u00B7 %s" % (d["tier"], TIER_COUNT, d["tier_name"].upper()),
             9, MUTED, align=PP_ALIGN.RIGHT, spacing=2)
    add_text(s, 0.5, 1.35, 9, 1.5, d["skill_name"], 52, CREAM, bold=True)
    add_text(s, 0.5, 3.05, 8.6, 1.0, u"\u201C%s\u201D" % d["catch_line"], 22,
             COPPER, italic=True)
    progress_bar(s, 0.5, 4.7, 5.0, d["tier"], d["lesson_number"],
                 d["total_lessons"], True)
    region_label = "Americas" if d.get("region") == "america" else "AU / NZ"
    add_text(s, 5.5, 4.92, 4.0, 0.25,
             "%s  \u00B7  %s" % (region_label, d.get("date_label", "")),
             10, MUTED, align=PP_ALIGN.RIGHT)


def skill_slide(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    add_text(s, 0.5, 0.4, 6, 0.25, "THE SKILL", 10, COPPERD, bold=True, spacing=3)
    add_text(s, 0.5, 0.66, 9, 0.7, d["skill_name"], 30, DARK, bold=True)

    add_round(s, 0.5, 1.65, 4.4, 3.2, CARD, line_hex=CARD_LINE, line_pt=1.0, radius=0.06)
    add_text(s, 0.75, 1.9, 3.9, 0.25, "WHAT IT IS", 11, COPPER, bold=True, spacing=1.5)
    add_text(s, 0.75, 2.25, 3.9, 2.45, d.get("skill_what", ""), 13.5, BODY)

    add_round(s, 5.1, 1.65, 4.4, 3.2, DARK, radius=0.06)
    add_text(s, 5.35, 1.9, 3.9, 0.25, "WHY IT WORKS", 11, COPPER, bold=True, spacing=1.5)
    add_text(s, 5.35, 2.25, 3.9, 2.45, d.get("skill_why", ""), 13.5, CREAM)

    progress_bar(s, 0.5, 5.15, 4.0, d["tier"], d["lesson_number"], d["total_lessons"], False)


def call_slide(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    rc = d.get("real_call", {}) or {}
    add_text(s, 0.5, 0.4, 9, 0.25, "THIS WEEK'S REAL CALL", 10, COPPERD,
             bold=True, spacing=3)
    add_text(s, 0.5, 0.68, 9, 0.85, rc.get("what_happened", ""), 13, BODY, italic=True)

    add_text(s, 0.5, 1.7, 4.4, 0.25, "HOW IT USUALLY GOES", 10, STEEL, bold=True, spacing=1.5)
    add_text(s, 5.1, 1.7, 4.4, 0.25, "HOW IT LANDS BETTER", 10, SAGE, bold=True, spacing=1.5)

    by = 2.05
    for turn in (rc.get("before") or [])[:3]:
        by = bubble(s, 0.5, by, 4.4, turn.get("who", ""), turn.get("line", ""), STEEL)
    ay = 2.05
    for turn in (rc.get("after") or [])[:3]:
        ay = bubble(s, 5.1, ay, 4.4, turn.get("who", ""), turn.get("line", ""), SAGE)


def cheat_slide(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_text(s, 0.5, 0.45, 9, 0.25, "AT THE DESK THIS WEEK", 10, COPPER, bold=True, spacing=3)
    add_text(s, 0.5, 0.72, 9, 0.6, "Say it like this", 28, CREAM, bold=True)

    phrases = (d.get("cheat_phrases") or [])[:3]
    n = max(len(phrases), 1)
    card_w = (9.0 - 0.4 * (n - 1)) / n
    for i, p in enumerate(phrases):
        x = 0.5 + i * (card_w + 0.4)
        add_round(s, x, 1.6, card_w, 2.0, PANEL, line_hex=COPPER, line_pt=1.2, radius=0.08)
        add_text(s, x + 0.15, 1.72, 0.5, 0.45, str(i + 1), 22, COPPER, bold=True)
        add_text(s, x + 0.18, 2.2, card_w - 0.36, 1.3, u"\u201C%s\u201D" % p, 13,
                 CREAM, italic=True)

    add_round(s, 0.5, 3.95, 9.0, 0.85, COPPER, radius=0.08)
    add_text(s, 0.7, 4.07, 9, 0.2, "REMEMBER", 8, DARK, bold=True, spacing=2)
    add_text(s, 0.7, 4.28, 8.6, 0.45, u"\u201C%s\u201D" % d["catch_line"], 17, DARK,
             bold=True, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    progress_bar(s, 0.5, 5.0, 4.0, d["tier"], d["lesson_number"], d["total_lessons"], True)

    notes = ("FACILITATOR GUIDE (5 min):\n%s\n\nBLACK SWAN \u2014 go deeper:\n%s"
             % (d.get("facilitator_guide", ""), d.get("black_swan", "")))
    s.notes_slide.notes_text_frame.text = notes


def _normalise(payload, region):
    """Fill in any missing structural fields so rendering never crashes."""
    d = dict(payload or {})
    d["region"] = region
    d.setdefault("total_lessons", 24)
    lesson = int(d.get("lesson_number", 1) or 1)
    d["lesson_number"] = lesson
    tier, tier_name = derive_tier(lesson)
    d.setdefault("tier", tier)
    d.setdefault("tier_name", tier_name)
    d.setdefault("pass_number", 1)
    d.setdefault("skill_name", "Difficult Conversations")
    d.setdefault("catch_line", "")
    d.setdefault("real_call", {})
    d.setdefault("cheat_phrases", [])
    return d


def generate_region(payload, outdir, region):
    if not payload:
        return None
    d = _normalise(payload, region)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cover_slide(prs, d)
    skill_slide(prs, d)
    call_slide(prs, d)
    cheat_slide(prs, d)

    path = os.path.join(outdir, "negotiation_%s.pptx" % region)
    prs.save(path)
    return path


# Allow standalone testing: python generate_negotiation_deck.py sample.json out/
if __name__ == "__main__":
    import json, sys, tempfile
    data = json.load(open(sys.argv[1])) if len(sys.argv) > 1 else {}
    out = sys.argv[2] if len(sys.argv) > 2 else tempfile.mkdtemp()
    print(generate_region(data, out, data.get("region", "aunz")))
