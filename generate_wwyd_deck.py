#!/usr/bin/env python3
"""
generate_wwyd_deck.py

Builds "What Would You Do?" team-discussion decks for Travellers Autobarn.
Claude (upstream in Make) reads the week's tickets, selects the best teaching
scenario per region, anonymises it, and supplies:
  - scenario      : the anonymised situation, ending in a prompt to discuss
  - discussion_points : 2-3 questions for the Friday huddle
  - facilitator_answer : manager-only "good practice" steer -> speaker notes

Tickets are split into two decks by region/country:
  - America deck   -> region "america" / country "united_states"
  - AUS/NZ deck    -> region "aunz"    / country "australia" or "new_zealand"

Usage (CLI):
    python generate_wwyd_deck.py --input wwyd.json --outdir ./out

Usage (HTTP, via app.py routes):
    POST /generate-wwyd-deck/america
    POST /generate-wwyd-deck/aunz
    Body: {
      "date_label": "June 2026",
      "scenarios": [
        {"region": "america", "scenario": "...", "discussion_points": ["...","..."],
         "facilitator_answer": "...", "ticket_number": "...", "branch": "...", "date": "..."}
      ]
    }
"""

import argparse
import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Indigo / teal / plum "discussion" palettes — distinct from the praise and
# Day-1 decks. Muted-professional but a touch more engaging.
PALETTES = [
    {"bg": "3D348B", "band": "241C5C", "ink": "FFFFFF", "mark": "C3BBE8"},  # indigo
    {"bg": "2A6F7F", "band": "184952", "ink": "FFFFFF", "mark": "A9D2DA"},  # deep teal
    {"bg": "7A3B69", "band": "4A2240", "ink": "FFFFFF", "mark": "DCB6D0"},  # plum
    {"bg": "44577A", "band": "29384F", "ink": "FFFFFF", "mark": "B4C4DD"},  # blue slate
]

REGION_TITLES = {
    "america": "What Would You Do? — USA",
    "aunz": "What Would You Do? — Australia & New Zealand",
}


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _set_fill(shape, hex_str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


def _normalise(value):
    if not value:
        return ""
    return value.strip().lower().replace("_", " ").replace("-", " ")


def _region_for_country(country):
    c = _normalise(country)
    if c in ("united states", "usa", "us", "united states of america", "america"):
        return "america"
    if c in ("australia", "new zealand", "aus", "nz", "aotearoa"):
        return "aunz"
    return None


def _region_of(item):
    """Resolve an item's region from an explicit 'region' field or its country."""
    r = _normalise(item.get("region"))
    if r in ("america", "usa", "us"):
        return "america"
    if r in ("aunz", "au nz", "aus nz", "anz"):
        return "aunz"
    return _region_for_country(item.get("country"))


def _prettify_branch(branch):
    if not branch:
        return ""
    return " ".join(w.capitalize() for w in branch.strip().replace("_", " ").replace("-", " ").split())


def _clean_date(date):
    if not date:
        return ""
    date = date.strip()
    iso = date.replace("Z", "+00:00")
    try:
        dt = datetime.fromisoformat(iso)
        return dt.strftime("%-d %B %Y")
    except (ValueError, TypeError):
        pass
    import re
    return re.sub(r"\s+\d{1,2}:\d{2}\s*(am|pm)?\s*$", "", date, flags=re.IGNORECASE)


def _coerce_list(value):
    """Normalise scenarios that may arrive as a JSON string, single dict, or list."""
    if isinstance(value, str):
        try:
            value = json.loads(value)
        except (ValueError, TypeError):
            return []
    if isinstance(value, dict):
        value = [value]
    if not isinstance(value, (list, tuple)):
        return []
    out = []
    for item in value:
        if isinstance(item, str):
            try:
                item = json.loads(item)
            except (ValueError, TypeError):
                continue
        if isinstance(item, dict):
            out.append(item)
    return out


def _coerce_points(points):
    """Discussion points may arrive as a list or a single newline/semicolon string."""
    if not points:
        return []
    if isinstance(points, str):
        raw = points.replace(";", "\n").splitlines()
        return [p.strip(" -•\t") for p in raw if p.strip(" -•\t")]
    if isinstance(points, (list, tuple)):
        return [str(p).strip(" -•\t") for p in points if str(p).strip(" -•\t")]
    return []


def _add_slide(prs, item, palette, region):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, palette["bg"])
    bg.shadow.inherit = False

    # Big decorative "?" motif, top-left
    qm = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(2.4), Inches(2.4))
    tf = qm.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "?"
    r.font.size = Pt(170)
    r.font.bold = True
    r.font.color.rgb = _rgb(palette["mark"])

    # Header banner, top-right
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(0.55), Inches(7.5), Inches(0.95)
    )
    _set_fill(band, palette["band"])
    band.shadow.inherit = False
    btf = band.text_frame
    btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.2)
    btf.margin_right = Inches(0.2)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = REGION_TITLES.get(region, "What Would You Do?")
    br.font.size = Pt(19)
    br.font.bold = True
    br.font.name = "Calibri"
    br.font.color.rgb = _rgb(palette["ink"])

    # White card holding the scenario + prompt + discussion points
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.95), Inches(11.13), Inches(4.5)
    )
    _set_fill(card, "FFFFFF")
    card.shadow.inherit = False
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP
    ctf.margin_left = Inches(0.5)
    ctf.margin_right = Inches(0.5)
    ctf.margin_top = Inches(0.4)
    ctf.margin_bottom = Inches(0.3)

    scenario = str(item.get("scenario", "") or "").strip()
    prevention = _coerce_points(item.get("prevention_points"))
    handling = _coerce_points(item.get("handling_points"))
    # Backward-compat: fall back to a single discussion_points list
    fallback = _coerce_points(item.get("discussion_points"))

    # Scenario paragraph
    sp = ctf.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sr = sp.add_run()
    sr.text = scenario
    sr.font.size = Pt(15)
    sr.font.name = "Calibri"
    sr.font.color.rgb = _rgb("2B2D42")
    sp.space_after = Pt(12)

    def _section(heading, pts):
        hp = ctf.add_paragraph()
        hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run()
        hr.text = heading
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.name = "Calibri"
        hr.font.color.rgb = _rgb(palette["bg"])
        hp.space_after = Pt(4)
        for pt in pts:
            dp = ctf.add_paragraph()
            dp.alignment = PP_ALIGN.LEFT
            dr = dp.add_run()
            dr.text = "•  " + pt
            dr.font.size = Pt(13)
            dr.font.name = "Calibri"
            dr.font.color.rgb = _rgb("3D3D4E")
            dp.space_after = Pt(3)

    if prevention or handling:
        if prevention:
            _section("How could we prevent this?", prevention)
        if handling:
            # small gap before the second section
            gap = ctf.add_paragraph()
            gap.space_after = Pt(2)
            _section("How would you handle the customer?", handling)
    elif fallback:
        _section("What would you do?", fallback)

    # Footer: Ticket / Branch / Date for the manager's reference
    ticket = str(item.get("ticket_number", "") or "").strip()
    branch = _prettify_branch(str(item.get("branch", "") or ""))
    date = _clean_date(str(item.get("date", "") or ""))

    foot = slide.shapes.add_textbox(Inches(1.1), Inches(6.6), Inches(11.13), Inches(0.7))
    ftf = foot.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    parts = []
    if ticket:
        parts.append(("Ticket #", ticket))
    if branch:
        parts.append(("Branch", branch))
    if date:
        parts.append(("Date", date))
    for i, (label, value) in enumerate(parts):
        rl = fp.add_run()
        rl.text = f"{label}: "
        rl.font.size = Pt(14)
        rl.font.bold = True
        rl.font.name = "Calibri"
        rl.font.color.rgb = _rgb(palette["ink"])
        rv = fp.add_run()
        rv.text = value
        rv.font.size = Pt(14)
        rv.font.name = "Calibri"
        rv.font.color.rgb = _rgb(palette["ink"])
        if i < len(parts) - 1:
            rs = fp.add_run()
            rs.text = "      •      "
            rs.font.size = Pt(14)
            rs.font.bold = True
            rs.font.name = "Calibri"
            rs.font.color.rgb = _rgb(palette["ink"])

    # Manager-only "good practice" answer -> speaker notes
    answer = str(item.get("facilitator_answer", "") or "").strip()
    if answer:
        notes = slide.notes_slide.notes_text_frame
        notes.text = "FACILITATOR NOTES (not shown on slide)\n\nSuggested good practice:\n" + answer

    return slide


def _build_deck(items, region, date_label):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cbg = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(cbg, "241C5C")
    cbg.shadow.inherit = False

    title_box = cover.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.0))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = REGION_TITLES.get(region, "What Would You Do?")
    tr.font.size = Pt(44)
    tr.font.bold = True
    tr.font.name = "Calibri"
    tr.font.color.rgb = _rgb("FFFFFF")

    if date_label:
        sub = cover.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.8))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = date_label
        sr.font.size = Pt(24)
        sr.font.name = "Calibri"
        sr.font.color.rgb = _rgb("C3BBE8")

    for idx, item in enumerate(items):
        palette = PALETTES[idx % len(PALETTES)]
        _add_slide(prs, item, palette, region)

    return prs


def generate(payload, outdir):
    """Build both decks. Returns list of written paths."""
    return [p for region in ("america", "aunz")
            for p in [generate_region(payload, outdir, region)] if p]


def generate_region(payload, outdir, region):
    """Build a single region's WWYD deck, or None if no scenario matches."""
    os.makedirs(outdir, exist_ok=True)

    if isinstance(payload, str):
        try:
            payload = json.loads(payload)
        except (ValueError, TypeError):
            payload = {}

    date_label = payload.get("date_label") or datetime.now().strftime("%B %Y")
    # Accept "scenarios" (preferred) or fall back to "praise"/"items" for reuse
    raw = payload.get("scenarios")
    if raw is None:
        raw = payload.get("items", payload.get("praise", []))
    scenarios = _coerce_list(raw or [])

    items = [s for s in scenarios if _region_of(s) == region]
    if not items:
        return None

    prs = _build_deck(items, region, date_label)
    filenames = {"america": "wwyd_america.pptx", "aunz": "wwyd_aunz.pptx"}
    path = os.path.join(outdir, filenames[region])
    prs.save(path)
    return path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--outdir", default="./out")
    args = ap.parse_args()
    with open(args.input) as f:
        payload = json.load(f)
    written = generate(payload, args.outdir)
    if not written:
        print("No scenarios matched a known region. Nothing generated.")
    for p in written:
        print(f"Wrote {p}")


if __name__ == "__main__":
    main()
