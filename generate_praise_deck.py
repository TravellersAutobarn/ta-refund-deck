#!/usr/bin/env python3
"""
generate_praise_deck.py

Builds bright, fun "Wall of Praise" decks for Travellers Autobarn.
One slide per praised Zendesk ticket. Each slide shows the ticket number,
branch (pickup location) and date, with a large blank space left for staff
to copy/paste the customer's praise.

Tickets are split into two decks by Country:
  - America deck   -> Country == "United States"
  - AUS/NZ deck    -> Country in {"Australia", "New Zealand"}

Usage (CLI):
    python generate_praise_deck.py --input praise.json --outdir ./out

Usage (HTTP, via app.py route /generate-praise-deck):
    POST a JSON body shaped like:
    {
      "date_label": "June 2026",
      "praise": [
        {"ticket_number": "53112", "branch": "San Francisco",
         "date": "June 3, 2026", "country": "United States"},
        ...
      ]
    }
"""

import argparse
import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Bright, fun rotating palettes. Each slide picks the next one in the list so
# consecutive slides always differ. (bg, accent/banner, quote-mark colour)
# ---------------------------------------------------------------------------
PALETTES = [
    {"bg": "FF6B6B", "band": "FFE66D", "ink": "2B2D42", "mark": "FFFFFF"},  # coral / sunshine
    {"bg": "4ECDC4", "band": "FFD93D", "ink": "1A535C", "mark": "FFFFFF"},  # teal / gold
    {"bg": "6A4C93", "band": "FFCA3A", "ink": "FFFFFF", "mark": "FFD6FF"},  # grape / amber
    {"bg": "1982C4", "band": "8AC926", "ink": "FFFFFF", "mark": "CDEAFF"},  # ocean / lime
    {"bg": "FF924C", "band": "FFCA3A", "ink": "3D2C2E", "mark": "FFFFFF"},  # mango / honey
    {"bg": "06D6A0", "band": "FFD166", "ink": "073B4C", "mark": "FFFFFF"},  # mint / butter
    {"bg": "EF476F", "band": "FFD166", "ink": "FFFFFF", "mark": "FFE5EC"},  # raspberry / butter
    {"bg": "118AB2", "band": "06D6A0", "ink": "FFFFFF", "mark": "CFF5E7"},  # blue / emerald
]

REGION_TITLES = {
    "america": "USA Customer Praise",
    "aunz": "Australia & New Zealand Customer Praise",
}


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _set_fill(shape, hex_str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


def _normalise(value):
    """Lowercase, strip, and turn underscores/hyphens into spaces."""
    if not value:
        return ""
    return value.strip().lower().replace("_", " ").replace("-", " ")


def _region_for_country(country):
    c = _normalise(country)
    if c in ("united states", "usa", "us", "united states of america", "america"):
        return "america"
    if c in ("australia", "new zealand", "aus", "nz", "aotearoa"):
        return "aunz"
    return None


def _prettify_branch(branch):
    """los_angeles -> Los Angeles, new_zealand -> New Zealand, etc."""
    if not branch:
        return ""
    return " ".join(w.capitalize() for w in branch.strip().replace("_", " ").replace("-", " ").split())


def _clean_date(date):
    """Drop a trailing time portion like '10:45 AM' if present."""
    if not date:
        return ""
    import re
    return re.sub(r"\s+\d{1,2}:\d{2}\s*(am|pm)?\s*$", "", date.strip(), flags=re.IGNORECASE)


def _add_slide(prs, item, palette, region):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Full-bleed bright background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, palette["bg"])
    bg.shadow.inherit = False

    # Big decorative quotation mark, top-left
    qm = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(3), Inches(2.4))
    tf = qm.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\u201C"  # left double quote
    r.font.size = Pt(200)
    r.font.bold = True
    r.font.color.rgb = _rgb(palette["mark"])

    # Header banner pill, top-right
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(0.55), Inches(6.9), Inches(0.95)
    )
    _set_fill(band, palette["band"])
    band.shadow.inherit = False
    btf = band.text_frame
    btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.2)
    btf.margin_right = Inches(0.2)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = REGION_TITLES.get(region, "Customer Praise")
    br.font.size = Pt(20)
    br.font.bold = True
    br.font.name = "Calibri"
    br.font.color.rgb = _rgb(palette["ink"])

    # Large light card = the blank space for pasting the praise
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.1), Inches(11.13), Inches(3.5)
    )
    _set_fill(card, "FFFFFF")
    card.shadow.inherit = False
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.4)
    ctf.margin_right = Inches(0.4)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "[ Paste customer praise here ]"
    cr.font.size = Pt(18)
    cr.font.italic = True
    cr.font.name = "Calibri"
    cr.font.color.rgb = _rgb("B0B0B0")

    # Footer detail row: Ticket / Branch / Date
    ticket = str(item.get("ticket_number", "") or "").strip()
    branch = _prettify_branch(str(item.get("branch", "") or ""))
    date = _clean_date(str(item.get("date", "") or ""))

    foot = slide.shapes.add_textbox(Inches(1.1), Inches(6.0), Inches(11.13), Inches(1.0))
    ftf = foot.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER

    parts = []
    if ticket:
        parts.append(("Ticket #", ticket))
    if branch:
        parts.append(("Branch", branch))
    if date:
        parts.append(("Date", date))

    for i, (label, value) in enumerate(parts):
        rl = fp.add_run()
        rl.text = f"{label}: "
        rl.font.size = Pt(18)
        rl.font.bold = True
        rl.font.name = "Calibri"
        rl.font.color.rgb = _rgb(palette["ink"])

        rv = fp.add_run()
        rv.text = value
        rv.font.size = Pt(18)
        rv.font.name = "Calibri"
        rv.font.color.rgb = _rgb(palette["ink"])

        if i < len(parts) - 1:
            rs = fp.add_run()
            rs.text = "      \u2022      "
            rs.font.size = Pt(18)
            rs.font.bold = True
            rs.font.name = "Calibri"
            rs.font.color.rgb = _rgb(palette["ink"])

    return slide


def _build_deck(items, region, date_label):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cbg = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(cbg, "2B2D42")
    cbg.shadow.inherit = False

    title_box = cover.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.0))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = REGION_TITLES.get(region, "Customer Praise")
    tr.font.size = Pt(48)
    tr.font.bold = True
    tr.font.name = "Calibri"
    tr.font.color.rgb = _rgb("FFFFFF")

    if date_label:
        sub = cover.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.8))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = date_label
        sr.font.size = Pt(24)
        sr.font.name = "Calibri"
        sr.font.color.rgb = _rgb("FFD166")

    for idx, item in enumerate(items):
        palette = PALETTES[idx % len(PALETTES)]
        _add_slide(prs, item, palette, region)

    return prs


def generate(payload, outdir):
    """Build both decks from a payload dict. Returns list of written paths."""
    os.makedirs(outdir, exist_ok=True)
    date_label = payload.get("date_label") or datetime.now().strftime("%B %Y")
    praise = payload.get("praise", []) or []

    buckets = {"america": [], "aunz": []}
    for item in praise:
        region = _region_for_country(item.get("country"))
        if region:
            buckets[region].append(item)

    written = []
    filenames = {
        "america": "praise_america.pptx",
        "aunz": "praise_aunz.pptx",
    }
    for region, items in buckets.items():
        if not items:
            continue
        prs = _build_deck(items, region, date_label)
        path = os.path.join(outdir, filenames[region])
        prs.save(path)
        written.append(path)

    return written


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="Path to praise JSON file")
    ap.add_argument("--outdir", default="./out", help="Output directory")
    args = ap.parse_args()

    with open(args.input) as f:
        payload = json.load(f)

    written = generate(payload, args.outdir)
    if not written:
        print("No praise items matched a known country. Nothing generated.")
    for p in written:
        print(f"Wrote {p}")


if __name__ == "__main__":
    main()
