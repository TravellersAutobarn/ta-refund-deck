#!/usr/bin/env python3
"""
generate_day1_deck.py

Builds bright, fun "Day 1 Issues" decks for Travellers Autobarn.
One slide per praised Zendesk ticket. Each slide shows the ticket number,
branch (pickup location) and date, with a large blank space left for staff
to copy/paste the customer's praise.

Tickets are split into two decks by Country:
  - America deck   -> Country == "United States"
  - AUS/NZ deck    -> Country in {"Australia", "New Zealand"}

Usage (CLI):
    python generate_day1_deck.py --input praise.json --outdir ./out

Usage (HTTP, via app.py route /generate-praise-deck):
    POST a JSON body shaped like:
    {
      "date_label": "June 2026",
      "praise": [
        {"ticket_number": "53112", "branch": "San Francisco",
         "date": "June 3, 2026", "country": "United States"},
        ...
      ]
    }
"""

import argparse
import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Bright, fun rotating palettes. Each slide picks the next one in the list so
# consecutive slides always differ. (bg, accent/banner, quote-mark colour)
# ---------------------------------------------------------------------------
PALETTES = [
    {"bg": "3A506B", "band": "1C2541", "ink": "FFFFFF", "mark": "9DB4CC"},  # slate / navy
    {"bg": "395B64", "band": "1E3A40", "ink": "FFFFFF", "mark": "A5C9CA"},  # teal slate
    {"bg": "5C3A4D", "band": "33212C", "ink": "FFFFFF", "mark": "C7A3B5"},  # muted plum
    {"bg": "44576D", "band": "29384A", "ink": "FFFFFF", "mark": "AEBED0"},  # steel blue
    {"bg": "4A5043", "band": "2C3026", "ink": "FFFFFF", "mark": "BFC6B4"},  # olive slate
    {"bg": "5A4A42", "band": "342A25", "ink": "FFFFFF", "mark": "C9B8AE"},  # taupe
]

REGION_TITLES = {
    "america": "USA Day 1 Issues",
    "aunz": "Australia & New Zealand Day 1 Issues",
}


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _set_fill(shape, hex_str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


def _normalise(value):
    """Lowercase, strip, and turn underscores/hyphens into spaces."""
    if not value:
        return ""
    return value.strip().lower().replace("_", " ").replace("-", " ")


def _region_for_country(country):
    c = _normalise(country)
    if c in ("united states", "usa", "us", "united states of america", "america"):
        return "america"
    if c in ("australia", "new zealand", "aus", "nz", "aotearoa"):
        return "aunz"
    return None


def _prettify_branch(branch):
    """los_angeles -> Los Angeles, new_zealand -> New Zealand, etc."""
    if not branch:
        return ""
    return " ".join(w.capitalize() for w in branch.strip().replace("_", " ").replace("-", " ").split())


def _clean_date(date):
    """Format the date nicely. Handles ISO 8601 (2026-05-29T00:45:25.000Z)
    and strips a trailing time like '10:45 AM' from plain strings."""
    if not date:
        return ""
    date = date.strip()

    # Try ISO 8601 first
    iso = date.replace("Z", "+00:00")
    try:
        dt = datetime.fromisoformat(iso)
        return dt.strftime("%-d %B %Y")
    except (ValueError, TypeError):
        pass

    import re
    return re.sub(r"\s+\d{1,2}:\d{2}\s*(am|pm)?\s*$", "", date, flags=re.IGNORECASE)


def _add_slide(prs, item, palette, region):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Full-bleed bright background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, palette["bg"])
    bg.shadow.inherit = False

    # Big decorative quotation mark, top-left
    qm = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(3), Inches(2.4))
    tf = qm.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\u201C"  # left double quote
    r.font.size = Pt(200)
    r.font.bold = True
    r.font.color.rgb = _rgb(palette["mark"])

    # Header banner pill, top-right
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(0.55), Inches(6.9), Inches(0.95)
    )
    _set_fill(band, palette["band"])
    band.shadow.inherit = False
    btf = band.text_frame
    btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.2)
    btf.margin_right = Inches(0.2)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = REGION_TITLES.get(region, "Customer Praise")
    br.font.size = Pt(20)
    br.font.bold = True
    br.font.name = "Calibri"
    br.font.color.rgb = _rgb(palette["ink"])

    # Large light card = the blank space for pasting the praise
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.1), Inches(11.13), Inches(3.5)
    )
    _set_fill(card, "FFFFFF")
    card.shadow.inherit = False
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.4)
    ctf.margin_right = Inches(0.4)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "[ Paste issue details here ]"
    cr.font.size = Pt(18)
    cr.font.italic = True
    cr.font.name = "Calibri"
    cr.font.color.rgb = _rgb("B0B0B0")

    # Footer detail row: Ticket / Branch / Date
    ticket = str(item.get("ticket_number", "") or "").strip()
    branch = _prettify_branch(str(item.get("branch", "") or ""))
    date = _clean_date(str(item.get("date", "") or ""))

    foot = slide.shapes.add_textbox(Inches(1.1), Inches(6.0), Inches(11.13), Inches(1.0))
    ftf = foot.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER

    parts = []
    if ticket:
        parts.append(("Ticket #", ticket))
    if branch:
        parts.append(("Branch", branch))
    if date:
        parts.append(("Date", date))

    for i, (label, value) in enumerate(parts):
        rl = fp.add_run()
        rl.text = f"{label}: "
        rl.font.size = Pt(18)
        rl.font.bold = True
        rl.font.name = "Calibri"
        rl.font.color.rgb = _rgb(palette["ink"])

        rv = fp.add_run()
        rv.text = value
        rv.font.size = Pt(18)
        rv.font.name = "Calibri"
        rv.font.color.rgb = _rgb(palette["ink"])

        if i < len(parts) - 1:
            rs = fp.add_run()
            rs.text = "      \u2022      "
            rs.font.size = Pt(18)
            rs.font.bold = True
            rs.font.name = "Calibri"
            rs.font.color.rgb = _rgb(palette["ink"])

    return slide


def _build_deck(items, region, date_label):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cbg = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(cbg, "2B2D42")
    cbg.shadow.inherit = False

    title_box = cover.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.0))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = REGION_TITLES.get(region, "Customer Praise")
    tr.font.size = Pt(48)
    tr.font.bold = True
    tr.font.name = "Calibri"
    tr.font.color.rgb = _rgb("FFFFFF")

    if date_label:
        sub = cover.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.8))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = date_label
        sr.font.size = Pt(24)
        sr.font.name = "Calibri"
        sr.font.color.rgb = _rgb("9DB4CC")

    for idx, item in enumerate(items):
        palette = PALETTES[idx % len(PALETTES)]
        _add_slide(prs, item, palette, region)

    return prs


def _coerce_praise(praise):
    """
    Make.com sometimes serialises array items as JSON strings, or the whole
    array as a single string. Normalise whatever arrives into a list of dicts.
    """
    # Whole praise value arrived as a string -> try to parse it into a list
    if isinstance(praise, str):
        try:
            praise = json.loads(praise)
        except (ValueError, TypeError):
            return []

    # A single object instead of an array -> wrap it in a list
    if isinstance(praise, dict):
        praise = [praise]

    if not isinstance(praise, (list, tuple)):
        return []

    out = []
    for item in praise:
        # Each item arrived as a JSON string -> parse it into a dict
        if isinstance(item, str):
            try:
                item = json.loads(item)
            except (ValueError, TypeError):
                continue
        if isinstance(item, dict):
            out.append(item)
    return out


def generate(payload, outdir):
    """Build both decks from a payload dict. Returns list of written paths."""
    return [p for region in ("america", "aunz")
            for p in [generate_region(payload, outdir, region)] if p]


def generate_region(payload, outdir, region):
    """
    Build a single region's deck ('america' or 'aunz').
    Returns the written path, or None if no praise items match that region.
    """
    os.makedirs(outdir, exist_ok=True)

    # Payload itself may arrive as a JSON string
    if isinstance(payload, str):
        try:
            payload = json.loads(payload)
        except (ValueError, TypeError):
            payload = {}

    date_label = payload.get("date_label") or datetime.now().strftime("%B %Y")
    praise = _coerce_praise(payload.get("praise", []) or [])

    items = [item for item in praise if _region_for_country(item.get("country")) == region]
    if not items:
        return None

    prs = _build_deck(items, region, date_label)
    filenames = {"america": "day1_america.pptx", "aunz": "day1_aunz.pptx"}
    path = os.path.join(outdir, filenames[region])
    prs.save(path)
    return path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="Path to praise JSON file")
    ap.add_argument("--outdir", default="./out", help="Output directory")
    args = ap.parse_args()

    with open(args.input) as f:
        payload = json.load(f)

    written = generate(payload, args.outdir)
    if not written:
        print("No praise items matched a known country. Nothing generated.")
    for p in written:
        print(f"Wrote {p}")


if __name__ == "__main__":
    main()
