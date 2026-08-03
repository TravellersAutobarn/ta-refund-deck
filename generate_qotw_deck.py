#!/usr/bin/env python3
"""
generate_qotw_deck.py
Builds "Question of the Week" (Branch Prep Check) decks for Travellers Autobarn.

Claude (upstream in Make) reads the week's tickets, finds the single most
common customer issue per region, and turns it into FOUR knowledge questions
aimed at the branch staff who prepared the vehicle and put it on the road.
Presented live on Teams: questions go up, staff answer in the chat, then the
answers are revealed. It doubles as a weekly QA pass on the Retell KB.

Split by region/country:
  - America deck   -> region "america" / country "united_states"
  - AUS/NZ deck    -> region "aunz"    / country "australia" or "new_zealand"

Payload (v2 — quiz format):
    {
      "date_label": "August 2026",
      "questions": [
        {"region": "aunz",
         "topic": "Water pump fuse",
         "headline": "Water pump & gas stove: what did we check before this van went out?",
         "context_line": "7 customers rang us about this in the past week. ...",
         "staff_questions": ["...", "...", "...", "..."],
         "staff_answers": ["Fuse: ...", "No water: ...", "Heat pin: ...", "At handover: ..."],
         "talking_point": "...", "kb_topic": "...", "times_asked": 7,
         "kb_check": "..."}
      ]
    }

Payload (v1 — legacy single-question format) is still rendered correctly:
    {"region": "aunz", "question": "...", "answer": "step; step; step",
     "talking_point": "...", "kb_topic": "...", "times_asked": 7, "kb_check": "..."}

Usage (CLI):
    python generate_qotw_deck.py --input qotw.json --outdir ./out

Usage (HTTP, via app.py routes):
    POST /generate-qotw-deck/america
    POST /generate-qotw-deck/aunz
"""
import argparse
import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Teal / green / blue "knowledge" palettes — distinct from the other decks.
PALETTES = [
    {"bg": "0E7C7B", "band": "08504F", "ink": "FFFFFF", "mark": "9FD8D6"},  # teal
    {"bg": "2A6F4E", "band": "184631", "ink": "FFFFFF", "mark": "A7D3BC"},  # green
    {"bg": "1B6CA8", "band": "0E4368", "ink": "FFFFFF", "mark": "AFD4EC"},  # blue
    {"bg": "8A5A2B", "band": "5A3A1A", "ink": "FFFFFF", "mark": "E0C3A0"},  # amber
]

COVER_TITLES = {
    "america": "Question of the Week — USA",
    "aunz": "Question of the Week — Australia & New Zealand",
}
REGION_SHORT = {"america": "USA", "aunz": "AU & NZ"}
COVER_SUBTITLE = "For branch teams: vehicle prep & handover"
ANSWER_HEADER = "The four answers — and what changes at the yard"

GREY = "6B6B7B"
NAVY = "2B2D42"
SLATE = "3D3D4E"

TEAMS_RUN_SHEET = (
    "HOW TO RUN THIS ON TEAMS (3 minutes)\n"
    "\n"
    "1. Share this slide and read the questions out. Do not reveal the answers yet.\n"
    "2. Ask each branch to type ONE answer into the Teams chat — keeps everyone "
    "involved without talking over each other.\n"
    "3. Give it 60 seconds, then move to the answer slide.\n"
    "\n"
    "The point to make: every one of these calls came from a customer already on "
    "the road. Anything we check or demonstrate at the yard is a call we never take."
)


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _set_fill(shape, hex_str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


def _normalise(value):
    if not value:
        return ""
    return value.strip().lower().replace("_", " ").replace("-", " ")


def _region_for_country(country):
    c = _normalise(country)
    if c in ("united states", "usa", "us", "united states of america", "america"):
        return "america"
    if c in ("australia", "new zealand", "aus", "nz", "aotearoa"):
        return "aunz"
    return None


def _region_of(item):
    r = _normalise(item.get("region"))
    if r in ("america", "usa", "us"):
        return "america"
    if r in ("aunz", "au nz", "aus nz", "anz"):
        return "aunz"
    return _region_for_country(item.get("country"))


def _coerce_list(value):
    if isinstance(value, str):
        try:
            value = json.loads(value)
        except (ValueError, TypeError):
            return []
    if isinstance(value, dict):
        value = [value]
    if not isinstance(value, (list, tuple)):
        return []
    out = []
    for item in value:
        if isinstance(item, str):
            try:
                item = json.loads(item)
            except (ValueError, TypeError):
                continue
        if isinstance(item, dict):
            out.append(item)
    return out


def _coerce_steps(answer):
    """Answer may be a paragraph, a list of steps, or a newline/semicolon string."""
    if not answer:
        return []
    if isinstance(answer, (list, tuple)):
        return [str(a).strip(" -•\t") for a in answer if str(a).strip(" -•\t")]
    if isinstance(answer, str):
        if "\n" in answer or ";" in answer:
            raw = answer.replace(";", "\n").splitlines()
            return [a.strip(" -•\t") for a in raw if a.strip(" -•\t")]
        return [answer.strip()]
    return [str(answer)]


def _string_list(value):
    """Coerce a field that should be a list of strings (tolerates a JSON string)."""
    if isinstance(value, str):
        try:
            value = json.loads(value)
        except (ValueError, TypeError):
            return []
    if not isinstance(value, (list, tuple)):
        return []
    return [str(v).strip() for v in value if str(v).strip()]


def _split_label(text):
    """'Fuse: battery box next to the kitchen.' -> ('Fuse:', 'battery box ...')."""
    if ": " not in text:
        return "", text
    label, rest = text.split(": ", 1)
    if len(label) > 28 or not rest:
        return "", text
    return label + ":", rest


def _numbered(paragraph, colour_hex, size_pt, marl_in=0.42):
    """Auto-numbered list item with a hanging indent.

    Children of <a:pPr> must appear in schema order:
    lnSpc, spcBef, spcAft, buClr, buSzPts, buFont, buAutoNum, ..., defRPr.
    Set paragraph spacing BEFORE calling this, then insert ahead of defRPr.
    """
    pPr = paragraph._p.get_or_add_pPr()
    marl = int(marl_in * 914400)
    pPr.set("marL", str(marl))
    pPr.set("indent", str(-marl))

    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buSzPts", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)

    buClr = pPr.makeelement(qn("a:buClr"), {})
    buClr.append(pPr.makeelement(qn("a:srgbClr"), {"val": colour_hex}))
    buSzPts = pPr.makeelement(qn("a:buSzPts"), {"val": str(int(size_pt * 100))})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Calibri"})
    buAutoNum = pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})

    anchor = pPr.find(qn("a:defRPr"))
    for el in (buClr, buSzPts, buFont, buAutoNum):
        if anchor is not None:
            anchor.addprevious(el)
        else:
            pPr.append(el)


def _run(paragraph, text, size, colour, bold=False, italic=False):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Calibri"
    r.font.color.rgb = _rgb(colour)
    return r


def _add_lightbulb(slide, color):
    """Simple lightbulb motif from basic shapes, top-left."""
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(0.35), Inches(1.3), Inches(1.3))
    _set_fill(bulb, color)
    bulb.shadow.inherit = False
    base = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.55), Inches(0.5), Inches(0.35)
    )
    _set_fill(base, color)
    base.shadow.inherit = False


def _frame(prs, palette, region, banner_text):
    """Shared slide chrome: background, lightbulb, header banner. Returns the slide."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, palette["bg"])
    bg.shadow.inherit = False

    _add_lightbulb(slide, palette["mark"])

    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(0.55), Inches(7.9), Inches(0.95)
    )
    _set_fill(band, palette["band"])
    band.shadow.inherit = False
    btf = band.text_frame
    btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.2)
    btf.margin_right = Inches(0.2)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    _run(bp, banner_text, 19, palette["ink"], bold=True)
    return slide


def _card(slide):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.95), Inches(11.13), Inches(4.5)
    )
    _set_fill(card, "FFFFFF")
    card.shadow.inherit = False
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.7)
    ctf.margin_right = Inches(0.7)
    return ctf


def _footer(slide, item, palette):
    kb_topic = str(item.get("kb_topic", "") or "").strip()
    times = item.get("times_asked")
    foot_bits = []
    if kb_topic:
        foot_bits.append(("KB topic", kb_topic))
    if times not in (None, "", 0, "0"):
        foot_bits.append(("Customer calls this week", str(times)))
    if not foot_bits:
        return
    foot = slide.shapes.add_textbox(Inches(1.1), Inches(6.6), Inches(11.13), Inches(0.7))
    ftf = foot.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    for i, (label, value) in enumerate(foot_bits):
        _run(fp, f"{label}: ", 14, palette["ink"], bold=True)
        _run(fp, value, 14, palette["ink"])
        if i < len(foot_bits) - 1:
            _run(fp, "      •      ", 14, palette["ink"], bold=True)


def _add_question_slide(prs, item, palette, region):
    """Slide 1: pose the questions — staff answer in the Teams chat before the reveal."""
    slide = _frame(
        prs, palette, region,
        f"Question of the Week — Branch Prep Check ({REGION_SHORT.get(region, '')})",
    )
    ctf = _card(slide)

    headline = str(item.get("headline") or item.get("question") or "").strip()
    context_line = str(item.get("context_line") or "").strip()
    questions = _string_list(item.get("staff_questions"))

    hp = ctf.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    _run(hp, headline, 22 if questions else 26, palette["bg"], bold=True)

    if not questions:
        # Legacy v1 payload: single question, no quiz list.
        hp.space_after = Pt(18)
        pp = ctf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        _run(pp, context_line or "How would you explain this to a customer?",
             16, GREY, italic=True)
        _footer(slide, item, palette)
        slide.notes_slide.notes_text_frame.text = TEAMS_RUN_SHEET
        return slide

    if context_line:
        cp = ctf.add_paragraph()
        cp.alignment = PP_ALIGN.CENTER
        cp.space_before = Pt(14)
        _run(cp, context_line, 15, GREY, italic=True)

    for i, q in enumerate(questions):
        qp = ctf.add_paragraph()
        qp.alignment = PP_ALIGN.LEFT
        qp.space_before = Pt(15 if i == 0 else 8)
        qp.space_after = Pt(0)
        _run(qp, q, 16, NAVY)
        _numbered(qp, palette["bg"], 16)

    _footer(slide, item, palette)
    slide.notes_slide.notes_text_frame.text = TEAMS_RUN_SHEET
    return slide


def _add_answer_slide(prs, item, palette, region):
    """Slide 2: the reveal — the answers, then the line staff can say to a customer."""
    short = REGION_SHORT.get(region, "")
    answers = _string_list(item.get("staff_answers"))
    slide = _frame(
        prs, palette, region,
        f"The Answers — Branch Prep Check ({short})" if answers
        else f"The Answer — {'USA' if region == 'america' else 'Australia & New Zealand'}",
    )
    talking = str(item.get("talking_point", "") or "").strip()

    if answers:
        ctf = _card(slide)

        hp = ctf.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        _run(hp, ANSWER_HEADER, 19, palette["bg"], bold=True)

        for i, a in enumerate(answers):
            label, rest = _split_label(a)
            ap = ctf.add_paragraph()
            ap.alignment = PP_ALIGN.LEFT
            ap.space_before = Pt(11 if i == 0 else 7)
            ap.space_after = Pt(0)
            if label:
                _run(ap, label + " ", 14, palette["bg"], bold=True)
            _run(ap, rest, 14, NAVY)
            _numbered(ap, palette["bg"], 14)

        if talking:
            tp = ctf.add_paragraph()
            tp.alignment = PP_ALIGN.LEFT
            tp.space_before = Pt(12)
            _run(tp, "Say it simply:  ", 14, palette["bg"], bold=True)
            _run(tp, talking, 14, SLATE, italic=True)
    else:
        # Legacy v1 payload: restate the question, then the answer/process bullets.
        ctf = _card(slide)
        ctf.vertical_anchor = MSO_ANCHOR.TOP
        ctf.margin_left = Inches(0.5)
        ctf.margin_right = Inches(0.5)
        ctf.margin_top = Inches(0.35)

        question = str(item.get("question", "") or "").strip()
        steps = _coerce_steps(item.get("answer"))

        qp = ctf.paragraphs[0]
        qp.alignment = PP_ALIGN.LEFT
        _run(qp, "Q:  " + question, 15, palette["bg"], bold=True)
        qp.space_after = Pt(10)

        ap = ctf.add_paragraph()
        ap.alignment = PP_ALIGN.LEFT
        _run(ap, "The answer / process", 13, GREY, bold=True)
        ap.space_after = Pt(4)

        multi = len(steps) > 1
        for s in steps:
            bp2 = ctf.add_paragraph()
            bp2.alignment = PP_ALIGN.LEFT
            _run(bp2, ("•  " + s) if multi else s, 15, NAVY)
            bp2.space_after = Pt(4)

        if talking:
            tp = ctf.add_paragraph()
            tp.alignment = PP_ALIGN.LEFT
            tp.space_before = Pt(10)
            _run(tp, "Say it simply:  ", 14, palette["bg"], bold=True)
            _run(tp, talking, 14, SLATE, italic=True)

    _footer(slide, item, palette)

    # Manager QA reminder + any KB-check note -> speaker notes (on the answer slide)
    kb_check = str(item.get("kb_check", "") or "").strip()
    notes_parts = ["KB QA CHECK (not shown on slide) — manager only",
                   "\nVerify these answers match the official KB / Retell article before Friday.",
                   "If they differ, update whichever is wrong — that's the weekly QA win."]
    if kb_check:
        notes_parts.append("\nClaude's note:\n" + kb_check)
    slide.notes_slide.notes_text_frame.text = "\n".join(notes_parts)
    return slide


def _build_deck(items, region, date_label):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cbg = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(cbg, "08504F")
    cbg.shadow.inherit = False

    title_box = cover.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.0))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    _run(tp, COVER_TITLES.get(region, "Question of the Week"), 42, "FFFFFF", bold=True)

    sub = cover.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.8))
    stf = sub.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    subtitle = f"{date_label}  •  {COVER_SUBTITLE}" if date_label else COVER_SUBTITLE
    _run(sp, subtitle, 24, "9FD8D6")

    for idx, item in enumerate(items):
        palette = PALETTES[idx % len(PALETTES)]
        _add_question_slide(prs, item, palette, region)
        _add_answer_slide(prs, item, palette, region)
    return prs


def _backfill_legacy_fields(item):
    """Keep `question` / `answer` populated for anything downstream that reads them.

    Claude no longer emits these (it doubled the output and caused truncation) —
    they cost nothing to rebuild here.
    """
    if not item.get("question") and item.get("headline"):
        item["question"] = item["headline"]
    if not item.get("answer"):
        answers = _string_list(item.get("staff_answers"))
        if answers:
            item["answer"] = "; ".join(answers)
    return item


def generate(payload, outdir):
    return [p for region in ("america", "aunz")
            for p in [generate_region(payload, outdir, region)] if p]


def generate_region(payload, outdir, region):
    os.makedirs(outdir, exist_ok=True)
    if isinstance(payload, str):
        try:
            payload = json.loads(payload)
        except (ValueError, TypeError):
            payload = {}

    date_label = payload.get("date_label") or datetime.now().strftime("%B %Y")
    raw = payload.get("questions")
    if raw is None:
        raw = payload.get("items", [])
    questions = [_backfill_legacy_fields(q) for q in _coerce_list(raw or [])]
    items = [q for q in questions if _region_of(q) == region]
    if not items:
        return None

    prs = _build_deck(items, region, date_label)
    filenames = {"america": "qotw_america.pptx", "aunz": "qotw_aunz.pptx"}
    path = os.path.join(outdir, filenames[region])
    prs.save(path)
    return path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--outdir", default="./out")
    args = ap.parse_args()
    with open(args.input) as f:
        payload = json.load(f)
    written = generate(payload, args.outdir)
    if not written:
        print("No questions matched a known region. Nothing generated.")
    for p in written:
        print(f"Wrote {p}")


if __name__ == "__main__":
    main()
